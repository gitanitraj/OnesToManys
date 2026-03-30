#!/usr/bin/env python3
"""
Generate OnesToManys Project Introduction Presentation
A 10-slide PowerPoint presentation for students
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "OnesToManys Project"
    subtitle.text = "Build a Full-Stack 3-Tier Application\n\nA Week-Long Individual Project\nDatabase • REST API • Frontend"


def add_bullet_slide(prs, title_text, bullet_points):
    """Helper function to add a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = title_text

    tf = body.text_frame
    tf.clear()

    for i, (text, level) in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(18) if level == 0 else Pt(16)


def add_why_matters_slide(prs):
    """Slide 2: Why This Project Matters"""
    bullet_points = [
        ("You are not just writing code—you're practicing how real systems are built:", 0),
        ("Model data that reflects real-world relationships", 1),
        ("Build APIs that teams can consume", 1),
        ("Connect backend logic to user-facing interfaces", 1),
        ("Work across database, server, and frontend layers", 1),
        ("", 0),
        ("This project bridges theory and practical application development", 0),
    ]
    add_bullet_slide(prs, "Why This Project Matters", bullet_points)


def add_what_you_build_slide(prs):
    """Slide 3: What You Will Build"""
    bullet_points = [
        ("A Complete 3-Tier Application:", 0),
        ("Presentation Layer: Web user interface (Vanilla JS + React)", 1),
        ("Business Logic Layer: REST API service", 1),
        ("Data Layer: Relational database with proper schema", 1),
        ("", 0),
        ("Core Capability:", 0),
        ("Full CRUD operations for master and detail entities", 1),
        ("One-to-many relationship endpoints and UI navigation", 1),
        ("Data import/export in multiple formats", 1),
    ]
    add_bullet_slide(prs, "What You Will Build", bullet_points)


def add_master_detail_slide(prs):
    """Slide 4: Core Concept - Master-Detail Relationships"""
    bullet_points = [
        ("Master-Detail (One-to-Many) Relationships:", 0),
        ("Master record: can exist independently", 1),
        ("Detail record: belongs to exactly one master", 1),
        ("", 0),
        ("Real-World Examples:", 0),
        ("Customer → Orders", 1),
        ("Course → Students", 1),
        ("Blog → Comments", 1),
        ("Playlist → Songs", 1),
        ("Project → Tasks", 1),
        ("", 0),
        ("Design Goal: Data integrity + intuitive navigation", 0),
    ]
    add_bullet_slide(prs, "Core Concept: Master-Detail Relationships", bullet_points)


def add_choose_domain_slide(prs):
    """Slide 5: Choose Your Own Domain"""
    bullet_points = [
        ("Pick ONE scenario and build it end-to-end yourself:", 0),
        ("", 0),
        ("Available Options:", 0),
        ("Customer-Orders management system", 1),
        ("Course-Students enrollment tracking", 1),
        ("Blog-Comments content management", 1),
        ("Playlist-Songs media organization", 1),
        ("Project-Tasks workflow management", 1),
        ("", 0),
        ("Rule of Thumb:", 0),
        ("Keep the model simple enough to finish in a week", 1),
        ("Make it interesting enough to stay motivated", 1),
    ]
    add_bullet_slide(prs, "Choose Your Own Domain", bullet_points)


def add_timeline_slide(prs):
    """Slide 6: Project Timeline (7 Days)"""
    bullet_points = [
        ("Days 1-2: Foundation", 0),
        ("Plan entities and relationships", 1),
        ("Write database schema SQL (DDL statements)", 1),
        ("Generate and load synthetic test data", 1),
        ("", 0),
        ("Days 3-4: Backend Development", 0),
        ("Implement REST API with CRUD endpoints", 1),
        ("Add one-to-many relationship routes", 1),
        ("Test with curl and Postman/Insomnia/Everest", 1),
        ("", 0),
        ("Days 5-7: Frontend Integration", 0),
        ("Build Vanilla JavaScript client", 1),
        ("Build React application", 1),
        ("Implement CRUD operations and relationship views", 1),
    ]
    add_bullet_slide(prs, "Project Timeline (7 Days)", bullet_points)


def add_tech_stack_slide(prs):
    """Slide 7: Technical Stack Options"""
    bullet_points = [
        ("Java Track:", 0),
        ("Spring Boot framework", 1),
        ("Spring Data / JPA / Hibernate ORM", 1),
        ("Maven project management", 1),
        ("", 0),
        ("Python Track:", 0),
        ("Flask or FastAPI framework", 1),
        ("SQLAlchemy ORM", 1),
        ("pip-based package management", 1),
        ("", 0),
        ("Universal Tools:", 0),
        ("Git version control • curl • Postman/Insomnia/Everest", 1),
    ]
    add_bullet_slide(prs, "Technical Stack Options", bullet_points)


def add_deliverables_slide(prs):
    """Slide 8: Required Deliverables"""
    bullet_points = [
        ("By the end of the project, you must submit:", 0),
        ("", 0),
        ("Database schema SQL file (DDL statements)", 1),
        ("Synthetic test data SQL file", 1),
        ("REST API with full CRUD for both entities", 1),
        ("Relationship endpoints for one-to-many navigation", 1),
        ("Data import/export support (SQL and/or JSON)", 1),
        ("Vanilla JavaScript frontend application", 1),
        ("React frontend application", 1),
        ("", 0),
        ("All code should be clean, maintainable, and well-organized", 0),
    ]
    add_bullet_slide(prs, "Required Deliverables", bullet_points)


def add_evaluation_slide(prs):
    """Slide 9: How You Will Be Evaluated"""
    bullet_points = [
        ("Assessment Criteria:", 0),
        ("", 0),
        ("Correct relationship design and schema quality", 1),
        ("API completeness and proper HTTP conventions", 1),
        ("Code quality and maintainability", 1),
        ("Frontend usability for hierarchical data navigation", 1),
        ("End-to-end integration and reliability", 1),
        ("", 0),
        ("Success looks like:", 0),
        ("A working application where users can manage both master and detail records smoothly through an intuitive interface", 1),
    ]
    add_bullet_slide(prs, "How You Will Be Evaluated", bullet_points)


def add_getting_started_slide(prs):
    """Slide 10: Getting Started - First 24 Hours"""
    bullet_points = [
        ("Your Kickoff Checklist:", 0),
        ("", 0),
        ("1. Pick your domain (Customer-Orders, Course-Students, etc.)", 1),
        ("2. Draw master and detail entities with attributes", 1),
        ("3. Define primary keys and foreign key constraints", 1),
        ("4. Create schema SQL and load sample data", 1),
        ("5. Test your first GET endpoints with curl", 1),
        ("6. Set up a Kanban board with at least 12 starter cards", 1),
        ("", 0),
        ("Milestone Checkpoint:", 0),
        ("If your schema, first GET endpoints, and Kanban board are working by end of Day 1, you're on track!", 1),
    ]
    add_bullet_slide(prs, "Getting Started: First 24 Hours", bullet_points)


def main():
    """Generate the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add all 10 slides
    add_title_slide(prs)
    add_why_matters_slide(prs)
    add_what_you_build_slide(prs)
    add_master_detail_slide(prs)
    add_choose_domain_slide(prs)
    add_timeline_slide(prs)
    add_tech_stack_slide(prs)
    add_deliverables_slide(prs)
    add_evaluation_slide(prs)
    add_getting_started_slide(prs)

    # Save the presentation
    output_file = "OnesToManys-Project-Introduction.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
